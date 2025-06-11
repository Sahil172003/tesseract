from fastapi import APIRouter, HTTPException, UploadFile, File, Form
from pydantic import BaseModel
from typing import List, Dict, Any, Optional
import os
import json
import tempfile
import sqlite3
from pathlib import Path
import pandas as pd
import numpy as np
from io import BytesIO
import base64
import logging
from datetime import datetime
import traceback
import requests

# Document processing imports
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation
try:
    import easyocr
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    
import html2text
from bs4 import BeautifulSoup
import markdown

# LangChain imports
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_anthropic import ChatAnthropic
from langchain.chains.question_answering import load_qa_chain
from langchain_huggingface import HuggingFaceEmbeddings
from langchain.schema import Document

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api/claude-chat", tags=["claude-chat"])

# Pydantic Models
class DocumentUploadResponse(BaseModel):
    status: str
    filename: str
    file_format: str
    content_preview: str
    content_length: int
    session_id: str

class ChatMessage(BaseModel):
    role: str
    content: str
    timestamp: str

class ChatRequest(BaseModel):
    question: str
    session_id: Optional[str] = None  # Make session_id optional
    model: Optional[str] = "claude-3-5-sonnet-20241022"

class ChatResponse(BaseModel):
    status: str
    message: str
    session_id: Optional[str] = None
    history: List[ChatMessage]
    document_info: Optional[Dict[str, str]] = None

class ModelListResponse(BaseModel):
    status: str
    models: List[Dict[str, str]]

class DocumentProcessor:
    """Enhanced document processor supporting multiple file formats"""
    
    def __init__(self):
        self.supported_formats = {
            'pdf': self._process_pdf,
            'txt': self._process_text,
            'csv': self._process_csv,
            'xlsx': self._process_excel,
            'xls': self._process_excel,
            'docx': self._process_docx,
            'doc': self._process_docx,
            'pptx': self._process_pptx,
            'ppt': self._process_pptx,
            'json': self._process_json,
            'html': self._process_html,
            'htm': self._process_html,
            'md': self._process_markdown,
            'xml': self._process_xml,
            'jpg': self._process_image,
            'jpeg': self._process_image,
            'png': self._process_image,
            'bmp': self._process_image,
            'tiff': self._process_image,
        }
        self.ocr_reader = None
    
    def _init_ocr(self):
        """Initialize OCR reader lazily"""
        if OCR_AVAILABLE and self.ocr_reader is None:
            try:
                self.ocr_reader = easyocr.Reader(['en'])
            except Exception as e:
                logger.warning(f"Failed to initialize OCR: {e}")
                self.ocr_reader = None
    
    def get_supported_formats(self) -> List[str]:
        """Return list of supported file formats"""
        return list(self.supported_formats.keys())
    
    def process_file(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Process file based on its extension"""
        file_extension = Path(filename).suffix.lower().lstrip('.')
        
        if file_extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        try:
            processor = self.supported_formats[file_extension]
            result = processor(file_content, filename)
            
            return {
                'filename': filename,
                'format': file_extension,
                'content': result,
                'status': 'success'
            }
        except Exception as e:
            raise Exception(f"Error processing {file_extension} file: {str(e)}")
    
    def _process_pdf(self, file_content: bytes, filename: str) -> str:
        """Process PDF files"""
        pdf_file = BytesIO(file_content)
        pdf_reader = PdfReader(pdf_file)
        text = ""
        
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        
        return text.strip()
    
    def _process_text(self, file_content: bytes, filename: str) -> str:
        """Process plain text files"""
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        for encoding in encodings:
            try:
                return file_content.decode(encoding)
            except UnicodeDecodeError:
                continue
        raise Exception("Could not decode text file with any supported encoding")
    
    def _process_csv(self, file_content: bytes, filename: str) -> str:
        """Process CSV files"""
        csv_file = BytesIO(file_content)
        encodings = ['utf-8', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                csv_file.seek(0)
                df = pd.read_csv(csv_file, encoding=encoding)
                break
            except (UnicodeDecodeError, pd.errors.EmptyDataError):
                continue
        else:
            raise Exception("Could not read CSV file with any supported encoding")
        
        # Convert to readable text format
        text = f"Dataset: {filename}\n"
        text += f"Shape: {df.shape[0]} rows, {df.shape[1]} columns\n"
        text += f"Columns: {', '.join(df.columns.tolist())}\n\n"
        text += "Sample data:\n"
        text += df.head(10).to_string(index=False)
        text += "\n\nSummary statistics:\n"
        text += df.describe().to_string()
        
        return text
    
    def _process_excel(self, file_content: bytes, filename: str) -> str:
        """Process Excel files"""
        excel_file = BytesIO(file_content)
        
        try:
            excel_data = pd.read_excel(excel_file, sheet_name=None)
            text = f"Excel file: {filename}\n"
            
            for sheet_name, df in excel_data.items():
                text += f"\n--- Sheet: {sheet_name} ---\n"
                text += f"Shape: {df.shape[0]} rows, {df.shape[1]} columns\n"
                text += f"Columns: {', '.join(df.columns.astype(str).tolist())}\n"
                text += "Sample data:\n"
                text += df.head(5).to_string(index=False)
                text += "\n"
            
            return text
        except Exception as e:
            excel_file = BytesIO(file_content)
            df = pd.read_excel(excel_file)
            return self._dataframe_to_text(df, filename)
    
    def _process_docx(self, file_content: bytes, filename: str) -> str:
        """Process Word documents"""
        docx_file = BytesIO(file_content)
        doc = DocxDocument(docx_file)
        
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text.strip())
                text += " | ".join(row_text) + "\n"
        
        return text.strip()
    
    def _process_pptx(self, file_content: bytes, filename: str) -> str:
        """Process PowerPoint presentations"""
        pptx_file = BytesIO(file_content)
        prs = Presentation(pptx_file)
        
        text = f"PowerPoint Presentation: {filename}\n\n"
        
        for i, slide in enumerate(prs.slides, 1):
            text += f"--- Slide {i} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
            text += "\n"
        
        return text.strip()
    
    def _process_json(self, file_content: bytes, filename: str) -> str:
        """Process JSON files"""
        try:
            json_data = json.loads(file_content.decode('utf-8'))
            
            text = f"JSON file: {filename}\n\n"
            text += "Structure:\n"
            text += json.dumps(json_data, indent=2, ensure_ascii=False)
            
            return text
        except json.JSONDecodeError as e:
            raise Exception(f"Invalid JSON format: {str(e)}")
    
    def _process_html(self, file_content: bytes, filename: str) -> str:
        """Process HTML files"""
        html_content = file_content.decode('utf-8')
        
        h = html2text.HTML2Text()
        h.ignore_links = False
        h.ignore_images = False
        
        soup = BeautifulSoup(html_content, 'html.parser')
        
        for script in soup(["script", "style"]):
            script.decompose()
        
        text = h.handle(str(soup))
        return text.strip()
    
    def _process_markdown(self, file_content: bytes, filename: str) -> str:
        """Process Markdown files"""
        md_content = file_content.decode('utf-8')
        
        html = markdown.markdown(md_content)
        soup = BeautifulSoup(html, 'html.parser')
        text = soup.get_text()
        
        result = f"Markdown file: {filename}\n\n"
        result += "Original markdown:\n"
        result += md_content
        result += "\n\nConverted text:\n"
        result += text
        
        return result
    
    def _process_xml(self, file_content: bytes, filename: str) -> str:
        """Process XML files"""
        xml_content = file_content.decode('utf-8')
        soup = BeautifulSoup(xml_content, 'xml')
        
        text = f"XML file: {filename}\n\n"
        text += "Content:\n"
        text += soup.get_text(separator='\n', strip=True)
        
        return text
    
    def _process_image(self, file_content: bytes, filename: str) -> str:
        """Process image files using OCR"""
        if not OCR_AVAILABLE:
            return f"Image file: {filename}\n\nOCR processing not available. Please install easyocr and pillow."
        
        self._init_ocr()
        
        if self.ocr_reader is None:
            return f"Image file: {filename}\n\nOCR initialization failed."
        
        try:
            image = Image.open(BytesIO(file_content))
            
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            image_array = np.array(image)
            results = self.ocr_reader.readtext(image_array)
            
            text = f"Image file: {filename}\n\n"
            text += "Extracted text from image:\n"
            
            for (bbox, extracted_text, confidence) in results:
                if confidence > 0.5:
                    text += f"{extracted_text}\n"
            
            return text.strip()
        except Exception as e:
            return f"Image file: {filename}\n\nError processing image: {str(e)}"
    
    def _dataframe_to_text(self, df: pd.DataFrame, filename: str) -> str:
        """Convert DataFrame to readable text"""
        text = f"Dataset: {filename}\n"
        text += f"Shape: {df.shape[0]} rows, {df.shape[1]} columns\n"
        text += f"Columns: {', '.join(df.columns.astype(str).tolist())}\n\n"
        text += "Sample data:\n"
        text += df.head(10).to_string(index=False)
        
        numeric_cols = df.select_dtypes(include=[np.number]).columns
        if len(numeric_cols) > 0:
            text += "\n\nSummary statistics:\n"
            text += df[numeric_cols].describe().to_string()
        
        return text

class ClaudeChatManager:
    """Claude-based document chat manager with session handling"""
    
    def __init__(self):
        self.processor = DocumentProcessor()
        self.sessions: Dict[str, Dict[str, Any]] = {}
        self.general_chat_history: List[Dict[str, Any]] = []  # For general chat without documents
        self.embeddings = None
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        self._init_embeddings()
    
    def _init_embeddings(self):
        """Initialize embeddings"""
        try:
            self.embeddings = HuggingFaceEmbeddings()
        except Exception as e:
            logger.warning(f"Failed to initialize embeddings: {e}")
            self.embeddings = None
    
    def create_session(self, file_content: bytes, filename: str) -> str:
        """Create a new chat session with uploaded document"""
        try:
            # Process document
            processed_doc = self.processor.process_file(file_content, filename)
            
            # Generate session ID
            session_id = f"claude_session_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{len(self.sessions)}"
            
            # Create vector store if embeddings available
            vector_store = None
            if self.embeddings:
                try:
                    chunks = self.text_splitter.split_text(processed_doc['content'])
                    documents = [Document(page_content=chunk) for chunk in chunks]
                    vector_store = FAISS.from_documents(documents, self.embeddings)
                except Exception as e:
                    logger.warning(f"Failed to create vector store: {e}")
            
            # Store session data
            self.sessions[session_id] = {
                'filename': filename,
                'content': processed_doc['content'],
                'format': processed_doc['format'],
                'vector_store': vector_store,
                'chat_history': [],
                'created_at': datetime.now().isoformat()
            }
            
            return session_id
            
        except Exception as e:
            logger.error(f"Error creating session: {e}")
            raise
    
    def get_session(self, session_id: str) -> Optional[Dict[str, Any]]:
        """Get session data"""
        return self.sessions.get(session_id)
    
    def add_message_to_history(self, session_id: Optional[str], role: str, content: str):
        """Add message to chat history (general or session-specific)"""
        message = {
            'role': role,
            'content': content,
            'timestamp': datetime.now().isoformat()
        }
        
        if session_id and session_id in self.sessions:
            self.sessions[session_id]['chat_history'].append(message)
        else:
            self.general_chat_history.append(message)
    
    def get_chat_history(self, session_id: Optional[str]) -> List[Dict[str, Any]]:
        """Get chat history (general or session-specific)"""
        if session_id and session_id in self.sessions:
            return self.sessions[session_id]['chat_history']
        else:
            return self.general_chat_history
    
    def answer_question(self, session_id: Optional[str], question: str, model: str = "claude-3-5-sonnet-20241022") -> str:
        """Answer question about document using Claude or general chat"""
        try:
            # Add user question to history
            self.add_message_to_history(session_id, "user", question)
            
            # Get Anthropic API key
            anthropic_api_key = os.getenv("ANTHROPIC_API_KEY")
            if not anthropic_api_key:
                fallback_response = self._get_fallback_response(session_id, question)
                self.add_message_to_history(session_id, "assistant", fallback_response)
                return fallback_response
            
            # Prepare context based on whether we have a document session
            if session_id and session_id in self.sessions:
                session = self.sessions[session_id]
                # Use vector store if available
                if session['vector_store'] and self.embeddings:
                    try:
                        docs = session['vector_store'].similarity_search(question, k=3)
                        context = "\n\n".join([doc.page_content for doc in docs])
                    except Exception as e:
                        logger.warning(f"Vector search failed: {e}")
                        context = session['content'][:8000]  # Claude can handle more context
                else:
                    context = session['content'][:8000]  # Use first 8000 chars for Claude
                
                # Format prompt with document context
                prompt = f"""Based on the following document content, please answer the user's question accurately and concisely.

Document: {session['filename']}
Content:
{context}

Question: {question}

Please provide a clear and informative answer based only on the information available in the document. If the information is not available in the document, please state that clearly."""
            else:
                # General chat without document context
                prompt = f"""You are Claude, a helpful AI assistant created by Anthropic. Please answer the following question in a clear, informative, and conversational way.

Question: {question}

Please provide a helpful and accurate response."""
            
            # Create Claude LLM instance
            llm = ChatAnthropic(
                model=model,
                anthropic_api_key=anthropic_api_key,
                temperature=0.1,
                max_tokens=4000
            )
            
            # Get response from Claude
            response = llm.invoke(prompt)
            answer = response.content if hasattr(response, 'content') else str(response)
            
            # Add AI response to history
            self.add_message_to_history(session_id, "assistant", answer)
            
            return answer
            
        except Exception as e:
            logger.error(f"Error answering question: {e}")
            fallback_response = f"I apologize, but I encountered an error while processing your question: {str(e)}. Please try rephrasing your question."
            self.add_message_to_history(session_id, "assistant", fallback_response)
            return fallback_response
    
    def _get_fallback_response(self, session_id: Optional[str], question: str) -> str:
        """Provide fallback response when Claude is unavailable"""
        if session_id and session_id in self.sessions:
            session = self.sessions[session_id]
            content = session['content']
            content_lower = content.lower()
            question_lower = question.lower()
            
            if any(word in question_lower for word in ['what', 'describe', 'explain']):
                # Extract first few sentences that might be relevant
                sentences = content.split('.')[:3]
                return f"Based on the document, here's what I found: {'. '.join(sentences)}."
            
            elif any(word in question_lower for word in ['how many', 'count', 'number']):
                return f"I can see numerical information in {session['filename']}, but I need Claude API access to provide accurate analysis."
            
            elif any(word in question_lower for word in ['summary', 'summarize', 'overview']):
                # Return first paragraph or chunk
                paragraphs = content.split('\n\n')
                first_para = paragraphs[0] if paragraphs else content[:500]
                return f"Here's a summary based on the document: {first_para}..."
            
            else:
                return f"I can see your question relates to {session['filename']}, but I need Claude API access to provide detailed analysis. Please ensure your Anthropic API key is configured."
        else:
            return "I can provide general assistance, but I need Claude API access to give you the best responses. Please ensure your Anthropic API key is configured."
    
    def clear_session(self, session_id: str) -> bool:
        """Clear a specific session"""
        if session_id in self.sessions:
            del self.sessions[session_id]
            return True
        return False
    
    def clear_general_chat(self):
        """Clear general chat history"""
        self.general_chat_history.clear()
    
    def get_all_sessions(self) -> List[Dict[str, Any]]:
        """Get all active sessions"""
        return [
            {
                'session_id': sid,
                'filename': data['filename'],
                'created_at': data['created_at'],
                'message_count': len(data['chat_history'])
            }
            for sid, data in self.sessions.items()
        ]

# Global chat manager instance
claude_chat_manager = ClaudeChatManager()

# Available Anthropic Claude models
ANTHROPIC_MODELS = [
    {"id": "claude-3-5-sonnet-20241022", "name": "Claude 3.5 Sonnet"},
    {"id": "claude-3-opus-20240229", "name": "Claude 3 Opus"},
    {"id": "claude-3-sonnet-20240229", "name": "Claude 3 Sonnet"},
    {"id": "claude-3-haiku-20240307", "name": "Claude 3 Haiku"},
]

# API Endpoints
@router.get("/models", response_model=ModelListResponse)
async def get_available_models():
    """Get list of available Anthropic Claude models"""
    return ModelListResponse(
        status="success",
        models=ANTHROPIC_MODELS
    )

@router.post("/upload", response_model=DocumentUploadResponse)
async def upload_document(file: UploadFile = File(...)):
    """Upload and process document for Claude chat"""
    try:
        # Validate file format
        supported_formats = claude_chat_manager.processor.get_supported_formats()
        file_extension = Path(file.filename).suffix.lower().lstrip('.')
        
        if file_extension not in supported_formats:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file format. Supported formats: {', '.join(supported_formats)}"
            )
        
        # Read file content
        file_content = await file.read()
        
        if len(file_content) == 0:
            raise HTTPException(status_code=400, detail="File is empty")
        
        logger.info(f"Processing file: {file.filename} ({len(file_content)} bytes)")
        
        # Create session
        session_id = claude_chat_manager.create_session(file_content, file.filename)
        session = claude_chat_manager.get_session(session_id)
        
        content_preview = session['content'][:500]
        if len(session['content']) > 500:
            content_preview += "..."
        
        return DocumentUploadResponse(
            status="success",
            filename=file.filename,
            file_format=session['format'],
            content_preview=content_preview,
            content_length=len(session['content']),
            session_id=session_id
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Upload error: {e}")
        raise HTTPException(status_code=500, detail=f"Upload failed: {str(e)}")

@router.post("/chat", response_model=ChatResponse)
async def chat_with_document(request: ChatRequest):
    """Chat with uploaded document using Claude or general chat"""
    try:
        logger.info(f"Processing Claude chat question: {request.question}")
        
        # Answer question (works with or without session_id)
        answer = claude_chat_manager.answer_question(
            request.session_id, 
            request.question, 
            request.model
        )
        
        # Get updated history
        history = [
            ChatMessage(
                role=msg['role'],
                content=msg['content'],
                timestamp=msg['timestamp']
            )
            for msg in claude_chat_manager.get_chat_history(request.session_id)
        ]
        
        # Get document info if session exists
        document_info = None
        if request.session_id:
            session = claude_chat_manager.get_session(request.session_id)
            if session:
                document_info = {
                    "filename": session['filename'],
                    "format": session['format']
                }
        
        return ChatResponse(
            status="success",
            message=answer,
            session_id=request.session_id,
            history=history,
            document_info=document_info
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Claude chat error: {e}")
        raise HTTPException(status_code=500, detail=f"Chat failed: {str(e)}")

@router.get("/sessions")
async def get_sessions():
    """Get all active Claude sessions"""
    try:
        sessions = claude_chat_manager.get_all_sessions()
        return {"status": "success", "sessions": sessions}
    except Exception as e:
        logger.error(f"Error getting sessions: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@router.delete("/sessions/{session_id}")
async def delete_session(session_id: str):
    """Delete a specific Claude session"""
    try:
        success = claude_chat_manager.clear_session(session_id)
        if success:
            return {"status": "success", "message": "Session deleted"}
        else:
            raise HTTPException(status_code=404, detail="Session not found")
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error deleting session: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@router.delete("/chat/clear")
async def clear_general_chat():
    """Clear general Claude chat history"""
    try:
        claude_chat_manager.clear_general_chat()
        return {"status": "success", "message": "General chat history cleared"}
    except Exception as e:
        logger.error(f"Error clearing general chat: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/test")
async def test_endpoint():
    """Test endpoint to verify Claude API functionality"""
    return {
        "status": "success",
        "message": "Claude document chat API is running",
        "timestamp": datetime.now().isoformat(),
        "supported_formats": claude_chat_manager.processor.get_supported_formats(),
        "active_sessions": len(claude_chat_manager.sessions)
    }

@router.get("/health")
async def health_check():
    """Health check endpoint for Claude API"""
    try:
        return {
            "status": "healthy",
            "timestamp": datetime.now().isoformat(),
            "active_sessions": len(claude_chat_manager.sessions),
            "ocr_available": OCR_AVAILABLE,
            "embeddings_available": claude_chat_manager.embeddings is not None,
            "anthropic_api_configured": bool(os.getenv("ANTHROPIC_API_KEY"))
        }
    except Exception as e:
        return {
            "status": "unhealthy",
            "error": str(e),
            "timestamp": datetime.now().isoformat()
        }